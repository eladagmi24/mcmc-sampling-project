import os
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)

TOTAL_SLIDES = 10


def dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def bar(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def text(slide, left, top, width, height, txt, size=18,
         color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = align


def bullets(slide, left, top, width, height, items, size=20,
            color=LIGHT_GRAY, spacing=Pt(12)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing


def slide_num(slide, num):
    text(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.3),
         '%d / %d' % (num, TOTAL_SLIDES), size=11, color=MEDIUM_GRAY,
         align=PP_ALIGN.RIGHT)


def title_bar(slide, title):
    text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         title, size=36, color=WHITE, bold=True)
    bar(slide, Inches(0.8), Inches(1.1), Inches(2.5), Pt(3), ACCENT_TEAL)


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(3.4), Inches(13.333), Pt(3), ACCENT_BLUE)
text(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
     "MCMC Sampling Methods for\nBayesian Regression",
     size=44, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(3.1), Inches(11), Inches(0.5),
     "Predicting Cloud Server Loads with Uncertainty Quantification",
     size=20, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
     "Advanced Methods in Machine Learning",
     size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
     "Elad Dagmi  |  Shaked Mizrahi",
     size=18, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(5.7), Inches(11), Inches(0.4),
     "July 2026", size=16, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem & Motivation
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_BLUE)
title_bar(s, "Problem & Motivation")

bullets(s, Inches(1.0), Inches(1.6), Inches(11), Inches(5), [
    "\u2022  Bayesian inference requires sampling from intractable posteriors",
    "\u2022  MCMC methods approximate distributions we cannot compute directly",
    "\u2022  Application: predict cloud server CPU load with uncertainty",
    "\u2022  Fundamental to modern ML: diffusion models, variational inference, RL",
    "\u2022  Bridges theory (Markov chains, detailed balance) and real-world practice",
    "\u2022  Goal: compare three MCMC samplers on a real-world regression task",
], size=22, spacing=Pt(20))

slide_num(s, 2)

# ============================================================
# SLIDE 3: Research Questions
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_PURPLE)
title_bar(s, "Research Questions")

questions = [
    ("Q1", "Convergence", "How fast does each method reach the target distribution?", ACCENT_BLUE),
    ("Q2", "Sampling Efficiency", "Which method gives the most independent samples per second?", ACCENT_TEAL),
    ("Q3", "Prediction Quality", "Does the choice of sampler affect RMSE and calibration?", ACCENT_ORANGE),
    ("Q4", "Robustness", "How sensitive is each method to its hyperparameters?", ACCENT_PURPLE),
]

for i, (qnum, qtitle, qdesc, qcolor) in enumerate(questions):
    y = 1.7 + i * 1.3
    text(s, Inches(1.0), Inches(y), Inches(1.0), Inches(0.4),
         qnum, size=22, color=qcolor, bold=True)
    text(s, Inches(2.0), Inches(y), Inches(4.0), Inches(0.4),
         qtitle, size=22, color=WHITE, bold=True)
    text(s, Inches(2.0), Inches(y + 0.4), Inches(10), Inches(0.4),
         qdesc, size=17, color=LIGHT_GRAY)

slide_num(s, 3)

# ============================================================
# SLIDE 4: Methods We Compare
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_BLUE)
title_bar(s, "Methods We Compare")

methods = [
    ("Metropolis-Hastings", ACCENT_BLUE, [
        "\u2022  General-purpose MCMC",
        "\u2022  Proposes from q(x'|x), accepts/rejects",
        "\u2022  Simple but random-walk behavior",
    ]),
    ("Gibbs Sampling", ACCENT_TEAL, [
        "\u2022  Samples each variable conditionally",
        "\u2022  100% acceptance rate",
        "\u2022  Requires known conditionals",
    ]),
    ("Hamiltonian Monte Carlo", ACCENT_ORANGE, [
        "\u2022  Uses gradient information",
        "\u2022  Efficient in high dimensions",
        "\u2022  Suppresses random-walk",
    ]),
]

for i, (name, color, items) in enumerate(methods):
    x = 0.8 + i * 4.1
    text(s, Inches(x), Inches(1.6), Inches(3.8), Inches(0.5),
         name, size=24, color=color, bold=True)
    bar(s, Inches(x), Inches(2.15), Inches(2.0), Pt(2), color)
    bullets(s, Inches(x), Inches(2.5), Inches(3.8), Inches(3),
            items, size=18, spacing=Pt(16))

slide_num(s, 4)

# ============================================================
# SLIDE 5: Comparison Parameters
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_TEAL)
title_bar(s, "Comparison Parameters")

params = [
    ("\u2022  Convergence rate", "trace plots, Gelman-Rubin R\u0302"),
    ("\u2022  Acceptance rate", "proposal efficiency"),
    ("\u2022  Effective Sample Size (ESS)", "independent samples after autocorrelation"),
    ("\u2022  Runtime efficiency", "ESS per second"),
    ("\u2022  Prediction RMSE", "test-set regression accuracy"),
    ("\u2022  Calibration", "95% credible interval coverage"),
]

for i, (param, desc) in enumerate(params):
    y = 1.7 + i * 0.85
    text(s, Inches(1.0), Inches(y), Inches(5), Inches(0.4),
         param, size=22, color=WHITE, bold=True)
    text(s, Inches(6.5), Inches(y), Inches(6), Inches(0.4),
         desc, size=18, color=MEDIUM_GRAY)

slide_num(s, 5)

# ============================================================
# SLIDE 6: Dataset
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_BLUE)
title_bar(s, "Dataset: Bitbrains Datacenter Traces")

bullets(s, Inches(1.0), Inches(1.6), Inches(5.5), Inches(4), [
    "\u2022  Source: GWA-T-12 (TU Delft / Kaggle)",
    "\u2022  1,750 VMs, ~30 days, 5-min intervals",
    "\u2022  Total size: ~284 MB (CSV)",
    "\u2022  Well-established in cloud research",
], size=20, spacing=Pt(20))

text(s, Inches(7.5), Inches(1.6), Inches(5), Inches(0.4),
     "Features per VM:", size=20, color=ACCENT_TEAL, bold=True)
bullets(s, Inches(7.5), Inches(2.2), Inches(5), Inches(3.5), [
    "\u2022  CPU usage (MHz) \u2014 target",
    "\u2022  CPU capacity, cores provisioned",
    "\u2022  Memory provisioned & usage",
    "\u2022  Disk read/write throughput",
    "\u2022  Network in/out throughput",
], size=18, spacing=Pt(14))

text(s, Inches(1.0), Inches(4.8), Inches(11), Inches(0.4),
     "Data Preparation:", size=20, color=ACCENT_ORANGE, bold=True)
bullets(s, Inches(1.0), Inches(5.4), Inches(11), Inches(1.5), [
    "\u2022  Lag features + rolling statistics + z-score normalization",
    "\u2022  70/30 temporal train/test split",
], size=18, spacing=Pt(14))

slide_num(s, 6)

# ============================================================
# SLIDE 7: Methodology
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_ORANGE)
title_bar(s, "Methodology")

# Simple flow: 6 steps as boxes with arrows
step_labels = [
    "Data Loading\n& EDA",
    "Feature\nEngineering",
    "Model\nDefinition",
    "MCMC\nSampling",
    "Diagnostics",
    "Evaluation",
]
step_colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE,
               ACCENT_ORANGE, ACCENT_BLUE, ACCENT_TEAL]

for i, (label, color) in enumerate(zip(step_labels, step_colors)):
    x = Inches(0.6 + i * 2.1)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8),
                              Inches(1.8), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for j, line in enumerate(label.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
    if i < len(step_labels) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  x + Inches(1.8), Inches(2.05),
                                  Inches(0.3), Inches(0.4))
        arr.fill.solid()
        arr.fill.fore_color.rgb = MEDIUM_GRAY
        arr.line.fill.background()

# Model definition
text(s, Inches(1.0), Inches(3.3), Inches(5), Inches(0.4),
     "Bayesian Linear Regression:", size=20, color=ACCENT_BLUE, bold=True)
bullets(s, Inches(1.0), Inches(3.9), Inches(5.5), Inches(2.5), [
    "\u2022  y | X, \u03b2, \u03c3\u00b2 ~ N(X\u03b2, \u03c3\u00b2I)",
    "\u2022  \u03b2 ~ N(0, \u03c4\u00b2I),  \u03c3\u00b2 ~ Inv-Gamma(a, b)",
    "\u2022  Conjugate priors \u2192 Gibbs conditionals",
], size=18, spacing=Pt(14))

text(s, Inches(7.0), Inches(3.3), Inches(5), Inches(0.4),
     "Tools:", size=20, color=ACCENT_TEAL, bold=True)
bullets(s, Inches(7.0), Inches(3.9), Inches(5.5), Inches(2.5), [
    "\u2022  Python (NumPy, SciPy, Matplotlib)",
    "\u2022  Custom MH, Gibbs, HMC implementations",
    "\u2022  ArviZ for MCMC diagnostics",
], size=18, spacing=Pt(14))

slide_num(s, 7)

# ============================================================
# SLIDE 8: Timeline
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_TEAL)
title_bar(s, "Timeline & Deliverables")

weeks = [
    ("Week 1", "Jul 1-7", "Data acquisition & EDA", ACCENT_BLUE),
    ("Week 2", "Jul 8-14", "Implement MH & Gibbs", ACCENT_TEAL),
    ("Week 3", "Jul 15-21", "Implement HMC, run experiments", ACCENT_ORANGE),
    ("Week 4", "Jul 22-28", "Diagnostics & comparison analysis", ACCENT_PURPLE),
    ("Week 5", "Jul 29 - Aug 4", "Write report & finalize notebook", ACCENT_BLUE),
    ("Week 6", "Aug 5-10", "Review, polish, submit", ACCENT_TEAL),
]

for i, (week, dates, task, color) in enumerate(weeks):
    y = 1.7 + i * 0.85
    bar(s, Inches(1.0), Inches(y), Pt(4), Inches(0.55), color)
    text(s, Inches(1.2), Inches(y + 0.05), Inches(1.5), Inches(0.4),
         week, size=20, color=color, bold=True)
    text(s, Inches(2.8), Inches(y + 0.05), Inches(2.5), Inches(0.4),
         dates, size=16, color=MEDIUM_GRAY)
    text(s, Inches(5.5), Inches(y + 0.05), Inches(7), Inches(0.4),
         task, size=20, color=LIGHT_GRAY)

text(s, Inches(1.0), Inches(7.0 - 0.7), Inches(11), Inches(0.4),
     "Deliverables:  Word report  +  Jupyter notebook  +  Comparative analysis",
     size=18, color=ACCENT_ORANGE, bold=True)

slide_num(s, 8)

# ============================================================
# SLIDE 9: References
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(0), Pt(5), Inches(7.5), ACCENT_BLUE)
title_bar(s, "References")

refs = [
    "Bishop, C.M. (2006). Pattern Recognition and Machine Learning. Ch. 11.",
    "Neal, R.M. (2011). MCMC Using Hamiltonian Dynamics. Handbook of MCMC.",
    "Gelman, A. et al. (2013). Bayesian Data Analysis, 3rd Ed.",
    "Robert & Casella (2004). Monte Carlo Statistical Methods.",
    "Betancourt, M. (2017). A Conceptual Introduction to HMC. arXiv:1701.02434.",
    "Hastings, W.K. (1970). Monte Carlo Sampling Methods Using Markov Chains.",
    "Shen, S. et al. (2015). Statistical Analysis of the Bitbrains Traces.",
]

bullets(s, Inches(1.0), Inches(1.6), Inches(11), Inches(5),
        refs, size=16, spacing=Pt(16))

slide_num(s, 9)

# ============================================================
# SLIDE 10: Thank You
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
bar(s, Inches(0), Inches(3.4), Inches(13.333), Pt(3), ACCENT_BLUE)
text(s, Inches(1), Inches(2.0), Inches(11), Inches(1),
     "Thank You", size=48, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
     "Questions?", size=28, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
     "Elad Dagmi  |  Shaked Mizrahi",
     size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save
DOCS_DIRECTORY = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'docs')
os.makedirs(DOCS_DIRECTORY, exist_ok=True)
out = os.path.join(DOCS_DIRECTORY, 'Sampling_Project_Proposal.pptx')
prs.save(out)
print('Saved: %s' % out)
print('Slides: %d' % len(prs.slides))
